import os
import re
import io
import requests
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from dotenv import load_dotenv
import google.generativeai as genai
import streamlit.components.v1 as components
import tempfile


# --- SESSION STATE DEFAULTS ---
if "step" not in st.session_state:
    st.session_state["step"] = 1
if "titles" not in st.session_state:
    st.session_state["titles"] = []
if "selected_title" not in st.session_state:
    st.session_state["selected_title"] = ""
if "slide_content" not in st.session_state:
    st.session_state["slide_content"] = []

# Load .env if present
load_dotenv()

# ---------- CONFIG ----------
DEFAULT_SLIDES = 5
AUDIENCE_PRESETS = ["Executive", "Technical", "Marketing", "Educational"]
MAX_WORDS_PER_SLIDE = 70
PEXELS_SEARCH_URL = "https://api.pexels.com/v1/search"

# ---------- STREAMLIT PAGE STYLE ----------
st.set_page_config(page_title="AI PPT Wizard", layout="wide")
st.markdown("""
    <style>
    .stApp { background: linear-gradient(rgba(6,8,12,0.92), rgba(6,8,12,0.92)); color: #e6eef8 !important; }
    .stButton>button { background-color:#2563EB !important; color: white !important; border: none !important; }
    .stSidebar { background-color:#071233 !important; color: #e6eef8 !important; }
    .stTextInput input, textarea { background-color:#0c1320 !important; color: #e6eef8 !important; border:1px solid #233554 !important; }
    .gradient-title {
      background: linear-gradient(90deg, #6ee7f9, #22d3ee, #60a5fa);
      -webkit-background-clip: text; -webkit-text-fill-color: transparent;
      font-weight: 900; letter-spacing: .3px; text-shadow: 0 0 .7px rgba(255,255,255,.15);
    }

    /* High-contrast “card” for sidebar sections */
    .sidebar-card {
      background:#0c1320; border:1px solid #233554; border-radius:12px;
      padding:14px; color:#e6eef8 !important; opacity:1 !important;
    }
    .sidebar-card h4 { margin:0 0 8px 0; font-weight:800; }
    .sidebar-card ol { margin:8px 0 0 18px; }
    .sidebar-card li { margin:4px 0; }
    .sidebar-card a { color:#8bd3ff; text-decoration:underline; }
    .badge-ok {
      display:inline-block; padding:2px 8px; border-radius:999px;
      background:#0f5132; color:#d1fae5; font-size:12px; margin-left:6px;
    }
    
    /* Labels for inputs, selects, sliders */
    .stTextInput label, .stSelectbox label, .stSlider label, .stTextArea label {
    color: #eaf2ff !important;          /* brighter */
    font-weight: 700 !important;         /* bold */
    letter-spacing: .2px;
    }

    /* Input & textarea fields */
    .stTextInput input, textarea, .stTextArea textarea {
    background-color: #0e1629 !important;
    color: #f7fbff !important;           /* bright text */
    border: 1px solid #3a4a6b !important;
    font-weight: 600 !important;
    }

    /* Placeholder text */
    .stTextInput input::placeholder,
    .stTextArea textarea::placeholder {
    color: #b9c7e6 !important;           /* lighter, still readable */
    opacity: 1 !important;
    }

    /* Select (dropdown) text container */
    div[data-baseweb="select"] {
    background-color: #0e1629 !important;
    color: #f7fbff !important;
    border: 1px solid #3a4a6b !important;
    font-weight: 600 !important;
    }

    /* Select menu items */
    div[role="listbox"] > div {
    color: #eaf2ff !important;
    }

    /* Slider numbers and ticks */
    .css-1siy2j7, .css-q8sbsg, .stSlider .st-c7,
    .stSlider [data-baseweb="slider"] {
    color: #eaf2ff !important;
    font-weight: 700 !important;
    }

    /* General small text (helper/descriptions) */
    .small, .stMarkdown p, .stCaption, .stAlert p {
    color: #dfe9ff !important;
    }
                
    /* ===== Force bright headers & labels (high priority) ===== */
    .stMarkdown h1, .stMarkdown h2, .stMarkdown h3,
    h1, h2, h3 {
    color: #eaf2ff !important;          /* bright */
    text-shadow: 0 0 1px rgba(255,255,255,.10);
    font-weight: 900 !important;
    }

    /* Step section header style */
    .step-header {
    color: #eaf2ff !important;          /* bright */
    font-weight: 900 !important;
    letter-spacing: .2px;
    border-left: 6px solid #22d3ee;
    padding-left: 12px;
    margin: 8px 0 16px 0;
    line-height: 1.25;
    }

    /* Labels above inputs */
    label, .stTextInput label, .stSelectbox label, .stSlider label, .stTextArea label {
    color: #eaf2ff !important;
    opacity: 1 !important;
    font-weight: 700 !important;
    }

    /* Inputs themselves (text is sometimes dim without this) */
    input, textarea, [data-baseweb="select"] * {
    color: #f7fbff !important;
    }

    /* Placeholder text (brighter) */
    input::placeholder, textarea::placeholder {
    color: #cbd8ff !important;
    opacity: 1 !important;
    }
    </style>
""", unsafe_allow_html=True)

# ---------- SIDEBAR ----------
with st.sidebar:
    st.title("⚡ AI PPT Wizard")

    # --- GEMINI KEY (persist + auto-apply with fallback to env) ---
    st.subheader("Gemini API Key")
    gemini_input = st.text_input(
        "Enter your Gemini API Key (optional)",
        type="password",
        help="Get a free key from [Google AI Studio](https://aistudio.google.com/app/apikey)"
    )
    if gemini_input:
        st.session_state["gemini_key"] = gemini_input.strip()

    G_API_KEY = st.session_state.get("gemini_key") or os.getenv("G_API_KEY")

    if not G_API_KEY:
        st.error("⚠ No Gemini API key found. Enter a key above (or set it in your environment) to generate slides.")
        st.stop()
    else:
        try:
            genai.configure(api_key=G_API_KEY)
            MODEL = genai.GenerativeModel("gemini-2.5-flash")
            st.markdown("**Gemini API key** <span class='badge-ok'>Active</span>", unsafe_allow_html=True)
        except Exception as e:
            st.error(f"Failed to initialize Gemini: {e}")
            st.stop()

    # --- PEXELS KEY (persist + fallback) ---
    st.subheader("Pexels API Key (Optional)")
    pexels_input = st.text_input(
        "Enter your Pexels API Key (optional)",
        type="password",
        help="Get a free key from [Pexels](https://www.pexels.com/api/)"
    )
    if pexels_input:
        st.session_state["pexels_key"] = pexels_input.strip()

    PEXELS_KEY = st.session_state.get("pexels_key") or os.getenv("PEXELS_API_KEY")

    if PEXELS_KEY:
        attach_images = True
        st.markdown("**Pexels images** <span class='badge-ok'>Enabled</span>", unsafe_allow_html=True)
    else:
        attach_images = False
        st.info("ℹ️ No Pexels API key provided. Slides will include image suggestions (keywords) only.")


    st.markdown("---")

    if st.button("📖 See Instructions", key="see_instructions"):
        st.markdown(
            """
            <div class="sidebar-card">
            <h4>✨ Flow</h4>
            <ol>
                <li><b>Enter Topic</b> & choose audience</li>
                <li><b>Generate titles</b> and pick one</li>
                <li><b>Edit</b> the outline/sections</li>
                <li><b>Preview & edit</b> bullets + notes</li>
                <li><b>Generate PPT</b> and download</li>
            </ol>
            <p>💡 Tip: Add a Pexels API key above to include real images.</p>
            </div>
            """,
            unsafe_allow_html=True
        )

    st.markdown("---")
    st.subheader("Disclaimer")
    st.markdown(
        """
        <div class="sidebar-card">
        <h4>⚠️ Important Notice</h4>
        <ul>
            <li>Slides are <b>AI-assisted</b> and may contain inaccuracies.</li>
            <li>Please <b>review and fact-check</b> before professional use.</li>
            <li>Your <b>API keys remain private</b> — used only locally.</li>
            <li>Images (if enabled) come from <b>Pexels</b> under their
            <a href="https://www.pexels.com/license/" target="_blank">license</a>.
            </li>
        </ul>
        </div>
        """,
        unsafe_allow_html=True
    )


# One-time popup
if 'popup_shown' not in st.session_state:
    st.session_state['popup_shown'] = True
    components.html("<script>alert('💡 This app uses a dark UI for best visibility.');</script>")

# ---------- HELPERS ----------
def safe_filename(s: str) -> str:
    return re.sub(r'[^a-zA-Z0-9_\-]+', '_', s).strip('_')[:80]

def word_count(text: str) -> int:
    return len(text.split())

def parse_lines_to_bullets_and_notes(generated_text: str):
    lines = [ln.strip() for ln in generated_text.splitlines() if ln.strip()]
    bullets, notes_lines = [], []
    for ln in lines:
        if re.match(r"^(-|•|\d+\.)\s*", ln):
            bullets.append(re.sub(r"^(-|•|\d+\.)\s*", "", ln).strip())
        else:
            notes_lines.append(ln)
    if not bullets and lines:
        bullets, notes_lines = lines[:4], lines[4:]
    return bullets[:6], " ".join(notes_lines).strip()

# ---------- PEXELS ----------
def fetch_image_url_safe(query, api_key):
    if not api_key:
        return None
    headers = {"Authorization": api_key}
    params = {"query": query, "per_page": 1}
    try:
        resp = requests.get(PEXELS_SEARCH_URL, headers=headers, params=params, timeout=10)
    except Exception as e:
        st.warning(f"Pexels request error: {e}")
        return None
    if resp.status_code == 200:
        try:
            data = resp.json()
            photos = data.get("photos") or []
            if photos:
                src = photos[0].get("src", {})
                return src.get("landscape") or src.get("original")
        except Exception as e:
            st.warning(f"Failed to parse Pexels response: {e}")
    elif resp.status_code in (401, 403, 429):
        st.warning("⚠️ Pexels API error — check your key or usage limits.")
    return None

def download_image_to_path(img_url, keyword):
    try:
        r = requests.get(img_url, timeout=15)
        r.raise_for_status()
        safe_kw = re.sub(r'[^a-z0-9]', '_', keyword.lower())[:40]
        fname = os.path.join(tempfile.gettempdir(), f"pexels_{safe_kw}.jpg")
        with open(fname, "wb") as f:
            f.write(r.content)
        return fname
    except Exception as e:
        st.warning(f"Failed to download image: {e}")
        return None


# ---------- LLM (Gemini) ----------
def generate_titles(subject, count=6):
    system = "You are an expert presentation author. Produce short, engaging presentation titles."
    prompt = f"Generate {count} concise titles (max 10 words each) for: \"{subject}\"."
    try:
        resp = MODEL.generate_content([system, prompt])
        text = resp.text or ""
        titles = [re.sub(r'^[\-\d\.\)\s]+', '', line).strip() for line in text.splitlines() if line.strip()]
        return titles[:count]
    except Exception as e:
        st.error(f"Error generating titles: {e}")
        return []

def generate_slide_text(ppt_title, section_title, audience, include_image_keyword=True):
    system = "You are an expert presentation writer. Use concise bullets and notes."
    prompt = (
        f"Create slide content for '{ppt_title}'.\n"
        f"Slide title: {section_title}\n"
        f"Audience: {audience}\n"
        "- 3 to 5 concise bullet points (<= 20 words each)\n"
        "- A short speaker note (1-2 sentences)\n"
    )
    if include_image_keyword:
        prompt += "- End with 'ImageKeyword: <2-4 word idea>'.\n"
    try:
        resp = MODEL.generate_content([system, prompt])
        text = resp.text or ""
        ik_match = re.search(r"ImageKeyword\s*:\s*(.+)$", text, flags=re.I | re.M)
        image_keyword = ik_match.group(1).strip() if ik_match else None
        if image_keyword:
            text = re.sub(r"ImageKeyword\s*:\s*.+$", "", text, flags=re.I | re.M).strip()
        bullets, notes = parse_lines_to_bullets_and_notes(text)
        return {"bullets": bullets, "notes": notes, "image_keyword": image_keyword}
    except Exception as e:
        st.error(f"Error generating slide: {e}")
        return {"bullets": ["(Generation failed)"], "notes": "", "image_keyword": None}

# ---------- PPT CREATION ----------
def create_pptx_bytes(ppt_title, slide_contents, attach_images=False):
    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = ppt_title
    try:
        slide.placeholders[1].text = "Generated by AI PPT Wizard"
    except Exception:
        pass

    for s in slide_contents:
        layout = prs.slide_layouts[1]
        slide_obj = prs.slides.add_slide(layout)
        slide_obj.shapes.title.text = s.get("slide_title", "")[:80]
        body = slide_obj.shapes.placeholders[1].text_frame
        body.clear()
        bullets = s.get("bullets", [])
        if bullets:
            body.paragraphs[0].text = bullets[0]
            for b in bullets[1:]:
                p = body.add_paragraph()
                p.text = b
        try:
            slide_obj.notes_slide.notes_text_frame.text = s.get("notes", "")
        except Exception:
            pass

        img_path = s.get("image_local_path")
        if attach_images and img_path and os.path.exists(img_path):
            try:
                slide_obj.shapes.add_picture(img_path, Inches(6.5), Inches(1.0), width=Inches(3.5))
            except Exception as e:
                st.warning(f"Could not add image: {e}")

    # Closing slide
    closing = prs.slides.add_slide(prs.slide_layouts[1])
    closing.shapes.title.text = "Conclusion & Next Steps"
    try:
        closing.shapes.placeholders[1].text_frame.text = "Summary and suggested next steps."
    except Exception:
        pass

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

# ---------- APP HEADER ----------
st.markdown("""
<style>
:root {
  --hero-grad: linear-gradient(90deg, #6ee7f9, #22d3ee, #60a5fa, #22d3ee, #6ee7f9);
}

/* fade-in + shimmer */
@keyframes shimmer {
  0% { background-position: 0% 50%; }
  100% { background-position: 200% 50%; }
}
@keyframes fadeScale {
  0% { opacity: 0; transform: translateY(4px) scale(0.98); }
  100% { opacity: 1; transform: translateY(0) scale(1); }
}

h1.hero-title{
  display:inline-block;               /* important for background-clip animation */
  background: var(--hero-grad);
  background-size: 200% 200%;
  -webkit-background-clip: text;
  -webkit-text-fill-color: transparent;
  animation: fadeScale .8s ease-out both, shimmer 3s linear infinite !important;
  font-weight: 900;
  letter-spacing: .3px;
  text-shadow: 0 0 1px rgba(255,255,255,.14);  /* subtle “brightness” */
}

</style>

<h1 class="hero-title">✨ AI PPT Wizard</h1>
""", unsafe_allow_html=True)

st.write("Guided flow: Topic → Titles → Outline → Edit → Generate PPT")



# ---------- STEP STATE HANDLERS ----------
def go_next(): st.session_state['step'] = min(5, st.session_state['step'] + 1)
def go_back(): st.session_state['step'] = max(1, st.session_state['step'] - 1)


# Step 1
if st.session_state['step'] == 1:
    st.markdown("<h2 class='step-header'>Step 1 — Topic & Purpose</h2>", unsafe_allow_html=True)
    topic = st.text_input("Enter the presentation topic", value=st.session_state.get('topic',''))
    audience = st.selectbox("Audience style", AUDIENCE_PRESETS, index=0)
    slides_count = st.slider("Desired number of slides", min_value=1, max_value=30, value=DEFAULT_SLIDES)
    if st.button("Generate Titles"):
        if not topic:
            st.error("Please enter a topic.")
        else:
            st.session_state['topic'] = topic
            st.session_state['audience'] = audience
            st.session_state['slides_count'] = slides_count
            with st.spinner("Generating titles ..."):
                titles = generate_titles(topic, count=8)
            if not titles:
                st.error("No titles generated. Try a different topic or check API key.")
            else:
                st.session_state['titles'] = titles
                go_next()
    st.markdown("---")
    st.write("Tip: Keep topics concise and descriptive (e.g., 'AI in Healthcare').")

# Step 2
if st.session_state['step'] == 2:
    st.markdown("<h2 class='step-header'>Step 2 — Pick & Edit Title</h2>", unsafe_allow_html=True)
    titles = st.session_state.get('titles', [])
    if not titles:
        st.warning("No titles yet — go back and generate.")
        if st.button("Back", key="back_step2_empty"):
            go_back()

    else:
        selected = st.radio("Choose a title", options=titles, index=0)
        custom_title = st.text_input("Or edit the chosen title", value=selected)
        cols = st.columns([1,1,1])
        with cols[0]:
            if st.button("Back", key="back_step2"):
                go_back()
        with cols[1]:
            if st.button("Regenerate Titles", key="regen_titles_step2"):
                with st.spinner("Regenerating..."):
                    st.session_state['titles'] = generate_titles(st.session_state.get('topic',''), count=8)
                st.rerun()
        with cols[2]:
            if st.button("Proceed", key="proceed_step2"):
                st.session_state['final_title'] = custom_title or selected
                topic = st.session_state.get('topic','Topic')
                default_sections = [
                    f"Introduction to {topic}",
                    f"Importance of {topic}",
                    f"Key technologies in {topic}",
                    "Case studies",
                    "Conclusion & Recommendations"
                ]
                st.session_state['sections'] = default_sections[:st.session_state.get('slides_count', DEFAULT_SLIDES)]
                go_next()


# Step 3
if st.session_state['step'] == 3:
    st.markdown("<h2 class='step-header'>Step 3 — Outline (Edit Sections)</h2>", unsafe_allow_html=True)
    sections = st.session_state.get('sections', [])
    if not sections:
        st.warning("No sections configured. Go back to select a title.")
        if st.button("Back"): go_back()
    else:
        edited = []
        for i, s in enumerate(sections):
            new_s = st.text_input(f"Section {i+1}", value=s, key=f"sec_{i}")
            edited.append(new_s)
        extra = st.text_input("Add an extra section (optional)")
        if extra:
            edited.append(extra)
        st.session_state['edited_sections'] = edited
        cols = st.columns([1,1])
        with cols[0]:
            if st.button("Back", key="back_step3"):
                go_back()
        with cols[1]:
            if st.button("Generate Slide Content", key="gen_slides_step3"):
                final_sections = edited[:st.session_state.get('slides_count', DEFAULT_SLIDES)]
                slide_contents = []
                with st.spinner("Generating slide content ..."):
                    for sec in final_sections:
                        slide = generate_slide_text(
                            st.session_state.get('final_title','Presentation'),
                            sec,
                            st.session_state.get('audience','Executive'),
                            include_image_keyword=True
                        )
                        slide['slide_title'] = sec
                        slide['image_local_path'] = None
                        # safe pexels flow
                        if attach_images and slide.get('image_keyword'):
                            img_url = fetch_image_url_safe(slide['image_keyword'], PEXELS_KEY)
                            if img_url:
                                local_path = download_image_to_path(img_url, slide['image_keyword'])
                                if local_path:
                                    slide['image_local_path'] = local_path
                        slide_contents.append(slide)
                st.session_state['slide_contents'] = slide_contents
                go_next()

# Step 4
if st.session_state['step'] == 4:
    st.markdown("<h2 class='step-header'>Step 4 — Preview & Edit Slides</h2>", unsafe_allow_html=True)
    slide_contents = st.session_state.get('slide_contents', [])
    if not slide_contents:
        st.warning("No slide content found. Generate slides first.")
        if st.button("Back"): go_back()
    else:
        for idx, s in enumerate(slide_contents):
            st.markdown(f"### Slide {idx+1}")
            title_in = st.text_input(f"Slide {idx+1} Title", value=s.get('slide_title',''), key=f"title_{idx}")
            bullets = s.get('bullets', [])
            new_bullets = []
            st.write("Bullets:")
            for j, b in enumerate(bullets):
                nb = st.text_input(f"Slide {idx+1} - Bullet {j+1}", value=b, key=f"bullet_{idx}_{j}")
                new_bullets.append(nb)
            if st.button(f"Add bullet to slide {idx+1}", key=f"add_b_{idx}"):
                updated = st.session_state['slide_contents']
                updated[idx]['bullets'] = new_bullets + ["(New bullet)"]
                st.session_state['slide_contents'] = updated
                st.rerun()
            notes = st.text_area(f"Speaker notes for slide {idx+1}", value=s.get('notes',''), key=f"notes_{idx}")
            image_kw = s.get('image_keyword')
            st.write(f"Image suggestion: **{image_kw if image_kw else 'No suggestion'}**")
            if s.get('image_local_path') and os.path.exists(s.get('image_local_path')):
                st.image(s.get('image_local_path'), width=320)
            if word_count(" ".join(new_bullets)) > MAX_WORDS_PER_SLIDE:
                st.warning(f"Slide {idx+1} exceeds recommended {MAX_WORDS_PER_SLIDE} words.")
            st.session_state['slide_contents'][idx]['slide_title'] = title_in
            st.session_state['slide_contents'][idx]['bullets'] = new_bullets
            st.session_state['slide_contents'][idx]['notes'] = notes

        cols = st.columns([1,1,1])
        with cols[0]:
            if st.button("Back", key="back_step4"):
                go_back()
        with cols[1]:
            if st.button("Regenerate this slide (not implemented)", key="regen_slide_step4"):
                st.info("Partial regeneration placeholder. You can edit bullets manually or regenerate whole flow.")
        with cols[2]:
            if st.button("Finalize & Create PPT", key="finalize_step4"):
                ppt_title = st.session_state.get('final_title','AI_Presentation')
                bio = create_pptx_bytes(ppt_title, st.session_state['slide_contents'], attach_images=attach_images)
                filename = safe_filename(ppt_title) + ".pptx"
                st.success("PPT created!")
                st.download_button(
                    "Download PPT",
                    data=bio,
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
                st.session_state['step'] = 5


# Step 5
if st.session_state['step'] == 5:
    st.markdown("<h2 class='step-header'>Step 5 — Finished</h2>", unsafe_allow_html=True)
    st.success("Presentation generated — check your downloads.")
    st.write("You can go back and tweak slides or start a new presentation.")
    if st.button("Start New Presentation", key="restart_step5"):
        keys = [
            'topic','audience','slides_count','titles',
            'final_title','sections','edited_sections','slide_contents'
        ]
        for k in keys:
            if k in st.session_state:
                del st.session_state[k]
        st.session_state['step'] = 1
        st.rerun()

